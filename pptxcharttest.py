from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
from pptx.util import Inches, Pt
import os
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml

# Establish current directory for downloads and saved files (Git ignores this)
current_dir = os.getcwd()
pptx_path = (f"{current_dir}/ola_report_template.pptx")
output_file = (f"{current_dir}/ola_report_outtest.pptx")
# print(pptx_path)

prs = Presentation(pptx_path)

# Pick the slide where the chart should be replaced
chart_slide = prs.slides[7]

# Find existing chart and get its coordinates
old_chart_shape = None
for shape in chart_slide.shapes:
    if shape.has_chart:
        old_chart_shape = shape
        break

if old_chart_shape:
    left = old_chart_shape.left
    top = old_chart_shape.top
    width = old_chart_shape.width
    height = old_chart_shape.height

    print(f"Old chart position: Left={left}, Top={top}, Width={width}, Height={height}")

    # Remove old chart
    chart_slide.shapes._spTree.remove(old_chart_shape._element)
else:
    # Fallback: default position if no chart found
    left, top, width, height = Inches(2), Inches(2), Inches(6), Inches(4.5)
    print("No existing chart found, using default position.")

# --- Remove old chart (optional) ---
for shape in list(chart_slide.shapes):
    if shape.has_chart:
        chart_slide.shapes._spTree.remove(shape._element)  # remove chart shape

# --- Prepare new chart data ---
chart_data = CategoryChartData()
chart_data.categories = ['1 YR NURI', '3 YR NURI', 'Right-sized on-demand', 'Like-for-like']
chart_data.add_series("Compute", (19.2, 21.4, 16.7, 25.3))
chart_data.add_series("Storage", (22.3, 24.5, 20.1, 27.8))
chart_data.add_series("Network", (22.3, 24.5, 20.1, 27.8))

# --- Add new chart ---
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = chart_slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data).chart

# Save updated presentation
prs.save(output_file)
print(f"Chart replaced successfully in {output_file}")


