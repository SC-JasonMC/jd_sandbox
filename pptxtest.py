
from pptx import Presentation
import os
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Establish current directory for downloads and saved files (Git ignores this)
current_dir = os.getcwd()
# pptx_path = "C:\\Users\\McleodJ\\OneDrive - SOFTCAT PLC\\Documents\\pytest\\Livv Housing Group MAS-64260 - AWS Migration Business Case.pptx"
# print(pptx_path)

input_directory = "C:\\Users\\McleodJ\\OneDrive - SOFTCAT PLC\\Documents\\pytest\\mgn\\OLD2"

pptx_path = [
    os.path.join(dirpath,f) 
    for (dirpath, dirnames, filenames) 
    in os.walk(input_directory) 
    for f in filenames
    if f.endswith("pptx")
][0]

def extract_text_from_shape(shape):
    """
    Recursively extract text from a shape.
    Handles grouped shapes and text-containing shapes.
    """
    text_content = ""

    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                # Join all paragraphs in the text frame
                text_content = "\n".join([p.text for p in shape.text_frame.paragraphs])
            elif shape.shape_type == 6:  # Group shape
                for sub_shape in shape.shapes:
                    text_content += extract_text_from_shape(sub_shape) + "\n"
        except Exception as e:
            text_content += f"[Error reading shape: {e}]\n"
    return text_content.strip()

# print(f"PPTX Path: {pptx_path}")

# print(pptx_path)

# files = [f for f in input_directory if os.path.isfile]
# print(files)

prs = Presentation(pptx_path)
slide = prs.slides[10]

read_pptx = extract_text_from_shape(shape=slide)

print(read_pptx)

# for shape in slide.shapes:
#     if shape.has_text_frame:
#         try:
#             if shape.has_text_frame:
#                 # Join all paragraphs in the text frame
#                 text_content = "\n".join([p.text for p in shape.text_frame.paragraphs])
#             elif shape.shape_type == 6:  # Group shape
#                 for sub_shape in shape.shapes:
#                     text_content += extract_text_from_shape(sub_shape) + "\n"
#         except Exception as e:
#             text_content += f"[Error reading shape: {e}]\n"
# #     if shape.has_table:
# #         table = shape.table
# #         for row in table.rows:
# #             print(row)
#             # for cell in row.cells:
#             #     if cell.text.strip():
#             #         for paragraph in cell.text_frame.paragraphs:
#             #             # Combine all runs text
#             #             full_text = "".join(run.text for run in paragraph.runs)
#             #             new_text = full_text
#             #             for old, new in replacements.items():
#             #                 new_text = new_text.replace(old, new)

#             #             if new_text != full_text:
#             #                 # Clear existing runs
#             #                 for run in paragraph.runs:
#             #                     run.text = ""
#             #                 # Add new text as a single run (preserves paragraph formatting)
#             #                 paragraph.runs[0].text = new_text


# # Load the PowerPoint file

# # Extract table data
# # for shape in slide.shapes:
# #     if shape.has_table: # Check if the shape is a table
# #         table = shape.table
# #         for row in table.rows:
# #             row_data = [cell.text for cell in row.cells]
# #             if row_data[0] == 'Storage':
# #                 storage_cost = row_data[1]
# #                 print(storage_cost)
# #             if row_data[0] == 'Network':
# #                 network_cost = row_data[1]
# #                 print(network_cost)
# #             if row_data[0] == 'Compute':
# #                 compute_cost = row_data[1]
# #                 print(compute_cost)


# # costs = {
# #     row_data[0]: row_data[1]
# #     for shape in slide.shapes if shape.has_table
# #     for row in shape.table.rows
# #     for row_data in [[cell.text.strip() for cell in row.cells]]
# #     if row_data[0] in ['Storage', 'Network', 'Compute']
# # }

# # print(costs['Storage'])